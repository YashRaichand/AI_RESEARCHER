import io
from typing import List, Dict, Any
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor


def _rgb(r, g, b):
    return RGBColor(r, g, b)


BG_COLOR = _rgb(10, 10, 30)
ACCENT = _rgb(127, 119, 221)
WHITE = _rgb(255, 255, 255)
GRAY = _rgb(180, 180, 200)


class SlideGenerator:
    def generate(self, papers_data: List[Dict[str, Any]], title: str) -> bytes:
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        self._add_title_slide(prs, title, f"{len(papers_data)} Research Papers Analyzed")

        self._add_overview_slide(prs, papers_data)

        for i, paper in enumerate(papers_data[:8]):
            self._add_paper_slide(prs, paper, i + 1)

        self._add_comparison_slide(prs, papers_data)
        self._add_references_slide(prs, papers_data)

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)
        return buf.read()

    def _set_bg(self, slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

    def _add_title(self, slide, text: str, left, top, width, height, size=28, bold=True, color=None):
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = text
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color or WHITE
        return txBox

    def _add_title_slide(self, prs, title: str, subtitle: str):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        self._set_bg(slide)
        self._add_title(slide, title, Inches(1), Inches(2.5), Inches(11), Inches(1.5), size=40, color=WHITE)
        self._add_title(slide, subtitle, Inches(1), Inches(4.2), Inches(11), Inches(1), size=20, bold=False, color=GRAY)
        self._add_title(slide, "Scientia AI — Research Assistant", Inches(1), Inches(6.5), Inches(11), Inches(0.5), size=14, bold=False, color=ACCENT)

    def _add_overview_slide(self, prs, papers_data: List[Dict]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide)
        self._add_title(slide, "Papers Overview", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), size=28, color=ACCENT)
        y = Inches(1.3)
        for i, paper in enumerate(papers_data[:6]):
            title_text = f"• {paper.get('title', 'Untitled')[:80]}"
            authors = ", ".join(str(a) for a in paper.get("authors", [])[:2])
            self._add_title(slide, title_text, Inches(0.5), y, Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
            y += Inches(0.38)
            if authors:
                self._add_title(slide, f"  {authors}", Inches(0.5), y, Inches(12), Inches(0.3), size=11, bold=False, color=GRAY)
                y += Inches(0.32)

    def _add_paper_slide(self, prs, paper: Dict, num: int):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide)
        title = paper.get("title", "Untitled")[:80]
        self._add_title(slide, f"Paper {num}: {title}", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), size=20, color=ACCENT)

        abstract = paper.get("abstract", "No abstract available.")[:600]
        self._add_title(slide, "Abstract:", Inches(0.5), Inches(1.3), Inches(12), Inches(0.4), size=14, bold=True, color=WHITE)
        self._add_title(slide, abstract, Inches(0.5), Inches(1.8), Inches(12), Inches(2.5), size=12, bold=False, color=GRAY)

        keywords = paper.get("keywords", [])[:8]
        if keywords:
            kw_text = "Keywords: " + " | ".join(keywords)
            self._add_title(slide, kw_text, Inches(0.5), Inches(4.5), Inches(12), Inches(0.5), size=12, bold=False, color=ACCENT)

        topic = paper.get("topic", "Other")
        quality = paper.get("quality_score")
        meta = f"Topic: {topic}"
        if quality:
            meta += f"  |  Quality Score: {quality:.0f}/100"
        self._add_title(slide, meta, Inches(0.5), Inches(5.2), Inches(12), Inches(0.4), size=12, bold=False, color=GRAY)

    def _add_comparison_slide(self, prs, papers_data: List[Dict]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide)
        self._add_title(slide, "Papers Comparison", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), size=28, color=ACCENT)

        rows = 1 + min(len(papers_data), 6)
        cols = 4
        table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12), Inches(5.5)).table
        headers = ["Title", "Topic", "Quality Score", "Status"]
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = header
            cell.fill.solid()
            cell.fill.fore_color.rgb = ACCENT
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.size = Pt(11)
                    run.font.bold = True
                    run.font.color.rgb = WHITE

        for row_idx, paper in enumerate(papers_data[:6], start=1):
            values = [
                paper.get("title", "Untitled")[:40],
                paper.get("topic", "Other"),
                f"{paper.get('quality_score', 0):.0f}/100" if paper.get("quality_score") else "N/A",
                paper.get("status", "completed"),
            ]
            for col_idx, val in enumerate(values):
                cell = table.cell(row_idx, col_idx)
                cell.text = val
                cell.fill.solid()
                cell.fill.fore_color.rgb = _rgb(20, 20, 45)
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.size = Pt(10)
                        run.font.color.rgb = WHITE

    def _add_references_slide(self, prs, papers_data: List[Dict]):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        self._set_bg(slide)
        self._add_title(slide, "References", Inches(0.5), Inches(0.3), Inches(12), Inches(0.8), size=28, color=ACCENT)
        y = Inches(1.3)
        for i, paper in enumerate(papers_data[:10]):
            title = paper.get("title", "Untitled")[:70]
            authors = ", ".join(str(a) for a in paper.get("authors", [])[:2])
            ref = f"[{i+1}] {authors + '. ' if authors else ''}{title}"
            self._add_title(slide, ref, Inches(0.5), y, Inches(12), Inches(0.45), size=11, bold=False, color=GRAY)
            y += Inches(0.48)
            if y > Inches(7.0):
                break
